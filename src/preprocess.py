import os
import csv
import warnings
warnings.filterwarnings("ignore")
from pptx import Presentation
import docx
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import NLTKTextSplitter
from langchain.docstore.document import Document
from dotenv import load_dotenv
load_dotenv()
import nltk
nltk.download('punkt')

def load_csv_documents(file_path):
    documents = []
    with open(file_path, mode='r', encoding='utf-8') as file:
        reader = csv.reader(file)
        header = next(reader)
        for row in reader:
            content = ' '.join(row)
            documents.append(Document(page_content=content, metadata={"source": file_path}))
    return documents

def load_pptx_documents(file_path):
    documents = []
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_texts.append(shape.text)
        content = "\n".join(slide_texts)
        documents.append(Document(page_content=content, metadata={"source": file_path}))
    return documents

def load_docx_documents(file_path):
    documents = []
    doc = docx.Document(file_path)
    content = []
    for para in doc.paragraphs:
        content.append(para.text)
    full_content = "\n".join(content)
    documents.append(Document(page_content=full_content, metadata={"source": file_path}))
    return documents

def load_documents(data_path):
    documents = []
    pdf_loader = PyPDFDirectoryLoader(data_path)
    pdf_documents = pdf_loader.load()
    documents.extend(pdf_documents)
    
    for file in os.listdir(data_path):
        if file.endswith(".csv"):
            csv_documents = load_csv_documents(os.path.join(data_path, file))
            documents.extend(csv_documents)
        elif file.endswith(".pptx"):
            ppt_documents = load_pptx_documents(os.path.join(data_path, file))
            documents.extend(ppt_documents)
        elif file.endswith(".docx"):
            docx_documents = load_docx_documents(os.path.join(data_path, file))
            documents.extend(docx_documents)
    return documents

def create_vector_db():
    data_path = os.getenv('DATA_PATH')
    db_path = os.getenv('DB_PATH')

    if not data_path or not db_path:
        raise ValueError("DATA_PATH or DB_PATH environment variables not set.")

    documents = load_documents(data_path)
    print(f"Processed {len(documents)} pages.")

    text_splitter = NLTKTextSplitter(chunk_size=1024, chunk_overlap=100)
    texts = text_splitter.split_documents(documents)
    print(f"Split into {len(texts)} chunks.")

    vector_store = Chroma.from_documents(
        documents=texts,
        embedding=HuggingFaceEmbeddings(),
        persist_directory=db_path
    )
    vector_store.persist()
    print(f"Vector database persisted at {db_path}.")

if __name__ == "__main__":
    create_vector_db()
